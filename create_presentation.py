#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for AI-Enabled ATS Application
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Style the title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    text_frame = body_shape.text_frame
    text_frame.clear()

    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(10)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(10)

def create_ats_presentation():
    """Create the complete ATS presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "AI-Enabled Applicant Tracking System",
        "Automated Resume Screening & Recruitment Management\nSpring Boot Application Project"
    )

    # Slide 2: Project Overview
    add_content_slide(
        prs,
        "Project Overview",
        [
            "🎯 Purpose: Streamline recruitment with AI-powered automation",
            "🤖 Core Feature: Intelligent resume screening algorithm",
            "👥 Target Users: HR teams, recruiters, and job candidates",
            "⚡ Key Benefit: Reduce manual screening time by 80%+",
            "🔐 Security: Role-based access control with Spring Security",
            "📊 Analytics: Comprehensive recruitment workflow tracking"
        ]
    )

    # Slide 3: Technology Stack
    add_content_slide(
        prs,
        "Technology Stack",
        [
            "Backend: Spring Boot 2.2.0 (Java 11)",
            "Security: Spring Security with Custom Authentication",
            "Database: PostgreSQL with JPA/Hibernate ORM",
            "Frontend: Thymeleaf Template Engine + HTML/CSS/JS",
            "Build Tool: Apache Maven",
            "AI Component: Custom resume analysis algorithm",
            "Code Quality: Lombok for clean code architecture"
        ]
    )

    # Slide 4: AI Resume Screening (Key Feature)
    add_content_slide(
        prs,
        "🤖 AI Resume Screening Algorithm",
        [
            "Automated Analysis: Instant evaluation upon application submission",
            "Skills Matching: Compares resume skills with job requirements",
            "Keyword Extraction: Identifies relevant technical terms",
            "Experience Scoring: Evaluates candidate background",
            "Match Score: Provides 0-100% compatibility rating",
            "Explainable AI: Shows matched keywords for transparency",
            "Location: AIResumeScreeningService.analyzeResume()"
        ]
    )

    # Slide 5: System Architecture
    add_content_slide(
        prs,
        "System Architecture (MVC Pattern)",
        [
            "Controllers: Role-based API endpoints",
            "  • AdminController - System administration",
            "  • RecruitmentController - Job & application management",
            "  • LoginController - Authentication handling",
            "Services: Business logic with AI integration",
            "  • AIResumeScreeningService - Core AI functionality",
            "  • RecruitmentService - Workflow orchestration",
            "Repositories: JPA data access with custom queries",
            "Models: Comprehensive recruitment domain entities"
        ]
    )

    # Slide 6: User Roles & Access Control
    add_two_column_slide(
        prs,
        "User Roles & Capabilities",
        [
            "👨‍💼 ADMIN",
            "• User management",
            "• System configuration",
            "• Full data access",
            "• Analytics dashboard",
            "",
            "👔 RECRUITER",
            "• Create job postings",
            "• Review applications",
            "• View AI scores",
            "• Schedule interviews"
        ],
        [
            "👤 CANDIDATE",
            "• Browse job openings",
            "• Submit applications",
            "• Upload resumes",
            "• Track status",
            "",
            "🔐 Security Features",
            "• BCrypt encryption",
            "• Domain-based roles",
            "• Session management",
            "• CSRF protection"
        ]
    )

    # Slide 7: Database Schema
    add_content_slide(
        prs,
        "Database Schema Design",
        [
            "📋 job_postings: Job details with required skills",
            "📝 applications: Submissions with AI scoring fields",
            "  • ai_score: Match percentage (0-100)",
            "  • ai_match_keywords: Extracted relevant terms",
            "👥 user_details: User profiles with role assignments",
            "📎 upload_files: Resume storage with text extraction",
            "🔗 Relationships: Optimized for recruitment workflow",
            "🗄️ PostgreSQL: Ensures data integrity and ACID compliance"
        ]
    )

    # Slide 8: Key Features
    add_two_column_slide(
        prs,
        "Key Features & Functionality",
        [
            "✅ Automated Features",
            "• AI resume screening",
            "• Auto score calculation",
            "• Email notifications",
            "• Status tracking",
            "",
            "📊 Recruitment Tools",
            "• Job posting management",
            "• Application dashboard",
            "• Interview scheduling",
            "• Candidate filtering"
        ],
        [
            "📁 File Management",
            "• Secure resume upload (3MB)",
            "• UUID-based naming",
            "• Text extraction",
            "• Multi-format support",
            "",
            "🔍 Advanced Search",
            "• Skill-based filtering",
            "• AI score sorting",
            "• Status-based views",
            "• Keyword search"
        ]
    )

    # Slide 9: AI Scoring Workflow
    add_content_slide(
        prs,
        "AI Resume Screening Workflow",
        [
            "1️⃣ Candidate uploads resume during application",
            "2️⃣ System extracts text from uploaded file",
            "3️⃣ AIResumeScreeningService.analyzeResume() triggered",
            "4️⃣ Algorithm processes:",
            "   • Skill matching against job requirements",
            "   • Keyword identification and extraction",
            "   • Experience level assessment",
            "5️⃣ Match score (0-100%) calculated and stored",
            "6️⃣ Recruiter views ranked applications with AI insights"
        ]
    )

    # Slide 10: Security Implementation
    add_content_slide(
        prs,
        "Security & Authentication",
        [
            "🔐 Spring Security Framework",
            "  • Custom authentication provider",
            "  • Role-based authorization",
            "  • Session management",
            "",
            "🛡️ Security Features",
            "  • BCrypt password encoding (default: Ats@ABC)",
            "  • Domain-based role assignment (@ats.com)",
            "  • Role hierarchy: ADMIN → RECRUITER → CANDIDATE",
            "  • CSRF token protection",
            "  • Secure file upload validation"
        ]
    )

    # Slide 11: Implementation Highlights
    add_content_slide(
        prs,
        "Implementation Highlights",
        [
            "✨ Clean Code: Lombok annotations reduce boilerplate",
            "🏗️ MVC Architecture: Separation of concerns",
            "🔄 Service Layer: Reusable business logic",
            "📦 Repository Pattern: Clean data access",
            "🎨 Thymeleaf Templates: Dynamic server-side rendering",
            "🧪 Testing: Functional test scripts provided",
            "📖 Documentation: Comprehensive CLAUDE.md guide",
            "🚀 Deployment: Docker-ready with health checks"
        ]
    )

    # Slide 12: Testing & Validation
    add_content_slide(
        prs,
        "Testing Strategy",
        [
            "🧪 Unit Testing: JUnit framework for service layer",
            "🔧 Functional Testing:",
            "  • test-ai-screening.sh - AI algorithm validation",
            "  • verify-implementation.sh - System verification",
            "  • RUN_APPLICATION.sh - Build and run checks",
            "📊 Manual Testing: Complete recruitment workflow",
            "🗄️ Database Verification: PostgreSQL query tests",
            "📁 File Upload Testing: Multi-format resume support",
            "🔐 Security Testing: Role-based access validation"
        ]
    )

    # Slide 13: Project Achievements
    add_content_slide(
        prs,
        "Project Achievements",
        [
            "✅ Fully functional recruitment management system",
            "✅ Working AI-powered resume screening algorithm",
            "✅ Complete role-based access control implementation",
            "✅ Secure file upload and management system",
            "✅ Responsive web interface with Thymeleaf",
            "✅ PostgreSQL database with optimized schema",
            "✅ Production-ready with health monitoring",
            "✅ Comprehensive documentation and test suite"
        ]
    )

    # Slide 14: Use Cases & Benefits
    add_two_column_slide(
        prs,
        "Real-World Impact",
        [
            "💼 For Organizations",
            "• 80%+ time savings",
            "• Reduced bias",
            "• Better candidates",
            "• Faster hiring",
            "• Cost effective",
            "",
            "📈 Scalability",
            "• High-volume hiring",
            "• Multiple recruiters",
            "• Concurrent users"
        ],
        [
            "👥 For HR Teams",
            "• Automated screening",
            "• AI-assisted decisions",
            "• Workflow tracking",
            "• Easy collaboration",
            "• Analytics insights",
            "",
            "🎯 For Candidates",
            "• Quick applications",
            "• Status transparency",
            "• Fair evaluation"
        ]
    )

    # Slide 15: Future Enhancements
    add_content_slide(
        prs,
        "Future Enhancements",
        [
            "🤖 Advanced AI: Machine learning models (NLP, GPT integration)",
            "📧 Email Integration: Automated candidate communication",
            "📱 Mobile App: Native iOS/Android applications",
            "🎥 Video Interviews: Integrated video screening",
            "📊 Advanced Analytics: Hiring metrics and dashboards",
            "🔗 Third-party Integration: LinkedIn, Indeed, job boards",
            "🌐 Multi-language Support: International recruitment",
            "☁️ Cloud Deployment: AWS/Azure scalable architecture"
        ]
    )

    # Slide 16: Demo Information
    add_content_slide(
        prs,
        "Live Demo Access",
        [
            "🌐 Application URL: http://localhost:8080/ats",
            "",
            "👤 Demo Credentials:",
            "  Admin: admin@ats.com / Admin@ABC",
            "  Test Candidates: Available with Ats@ABC password",
            "",
            "🔧 Setup Requirements:",
            "  • PostgreSQL database running",
            "  • Java 11+ installed",
            "  • Maven or IDE with Lombok plugin",
            "",
            "📦 Quick Start: Run GetreadyApplication.java from IDE"
        ]
    )

    # Slide 17: Technical Challenges
    add_content_slide(
        prs,
        "Technical Challenges Overcome",
        [
            "⚙️ Lombok Integration: Annotation processing setup",
            "🧠 AI Algorithm Design: Balancing accuracy and performance",
            "🔐 Security Configuration: Multi-role authentication flow",
            "📁 File Management: Secure handling of resume uploads",
            "🗄️ Database Optimization: Query performance tuning",
            "🎨 UI/UX Design: Intuitive recruitment workflow",
            "🧪 Testing Strategy: Comprehensive validation approach",
            "📖 Documentation: Clear guidelines for maintenance"
        ]
    )

    # Slide 18: Conclusion
    add_content_slide(
        prs,
        "Conclusion",
        [
            "✅ Successfully implemented AI-enabled ATS system",
            "🎯 Achieved goal of automated resume screening",
            "🏗️ Built scalable Spring Boot architecture",
            "🔐 Ensured security with role-based access control",
            "📊 Demonstrated practical AI application in HR",
            "🚀 Production-ready recruitment management solution",
            "",
            "💡 Key Takeaway: AI can significantly improve recruitment",
            "   efficiency while maintaining quality and fairness"
        ]
    )

    # Slide 19: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add centered title
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(2)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "Questions & Answers"

    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(54)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(0, 51, 102)
    paragraph.alignment = PP_ALIGN.CENTER

    # Add subtitle
    left = Inches(2)
    top = Inches(4)
    width = Inches(6)
    height = Inches(1)

    text_box2 = slide.shapes.add_textbox(left, top, width, height)
    text_frame2 = text_box2.text_frame
    text_frame2.text = "Thank you for your attention!"

    paragraph2 = text_frame2.paragraphs[0]
    paragraph2.font.size = Pt(28)
    paragraph2.font.color.rgb = RGBColor(68, 68, 68)
    paragraph2.alignment = PP_ALIGN.CENTER

    # Save the presentation
    filename = "ATS_Project_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Presentation created successfully: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return filename

if __name__ == "__main__":
    create_ats_presentation()
